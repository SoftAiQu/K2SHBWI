"""
PPTX Converter - Convert K2SHBWI files to PowerPoint presentations
"""

import io
from pathlib import Path
from typing import Dict
from PIL import Image

try:
    from pptx import Presentation as PPTXPresentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False
    PPTXPresentation = object

from .base_converter import BaseConverter


class PPTXConverter(BaseConverter):
    """Convert K2SHBWI to PowerPoint presentation"""
    
    @property
    def format_name(self) -> str:
        return 'PPTX'
    
    def _convert_impl(self, image: Image.Image, metadata: Dict, hotspots: list, output_path: str):
        """Convert to PPTX format"""
        
        if not HAS_PPTX:
            raise ImportError("python-pptx is required for PPTX conversion. Install with: pip install python-pptx")
        
        # Create presentation
        prs = PPTXPresentation()
        prs.slide_width = Inches(10)  # type: ignore
        prs.slide_height = Inches(7.5)  # type: ignore
        
        # Slide 1: Title slide
        self._add_title_slide(prs, metadata)
        
        # Slide 2: Main image with hotspots
        self._add_image_slide(prs, image, metadata, hotspots)
        
        # Slide 3: Hotspots details (if there are hotspots)
        if hotspots:
            self._add_hotspots_slide(prs, hotspots)
        
        # Slide 4: Metadata (if available)
        if metadata:
            self._add_metadata_slide(prs, metadata)
        
        # Save presentation
        prs.save(output_path)  # type: ignore
    
    def _add_title_slide(self, prs, metadata: Dict):
        """Add title slide to presentation"""
        
        # Use title slide layout
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Add background color
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(102, 126, 234)  # #667eea
        
        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = metadata.get('title', 'K2SHBWI Presentation')
        title_frame.paragraphs[0].font.size = Pt(54)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Add description
        desc_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(2))
        desc_frame = desc_box.text_frame
        desc_frame.word_wrap = True
        desc_frame.text = metadata.get('description', 'Interactive Image Presentation')
        desc_frame.paragraphs[0].font.size = Pt(28)
        desc_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        desc_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Add date/info at bottom
        info_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.8))
        info_frame = info_box.text_frame
        info_frame.text = "Generated from K2SHBWI Format"
        info_frame.paragraphs[0].font.size = Pt(14)
        info_frame.paragraphs[0].font.color.rgb = RGBColor(200, 200, 200)
        info_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    def _add_image_slide(self, prs, image: Image.Image, metadata: Dict, hotspots: list):
        """Add slide with main image"""
        
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
        title_frame = title_box.text_frame
        title_frame.text = "Main Image"
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(102, 126, 234)
        
        # Add image
        img_buffer = io.BytesIO()
        image.save(img_buffer, format='PNG')
        img_buffer.seek(0)
        
        # Calculate image dimensions
        max_width = Inches(8)
        max_height = Inches(5)
        img_width, img_height = image.size
        ratio = img_width / img_height
        
        if ratio > (8 / 5):  # Width is limiting
            display_width = max_width
            display_height = Inches(max_width.inches / ratio)
        else:  # Height is limiting
            display_height = max_height
            display_width = Inches(max_height.inches * ratio)
        
        left = Inches((10 - display_width.inches) / 2)
        top = Inches(1.2)
        
        slide.shapes.add_picture(img_buffer, left, top, width=display_width)
        
        # Add hotspot count
        if hotspots:
            info_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(0.5))
            info_frame = info_box.text_frame
            info_frame.text = f"Contains {len(hotspots)} interactive hotspots"
            info_frame.paragraphs[0].font.size = Pt(12)
            info_frame.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)
            info_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    def _add_hotspots_slide(self, prs, hotspots: list):
        """Add slide with hotspots details"""
        
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
        title_frame = title_box.text_frame
        title_frame.text = f"Interactive Hotspots ({len(hotspots)} total)"
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(102, 126, 234)
        
        # Add hotspots as bullet list
        content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(8.4), Inches(5.5))
        text_frame = content_box.text_frame
        text_frame.word_wrap = True
        
        for i, hotspot in enumerate(hotspots[:15], 1):  # Limit to 15 per slide
            title_hs = hotspot.get('data', {}).get('user_data', {}).get('title', f'Hotspot {i}')
            description_hs = hotspot.get('data', {}).get('user_data', {}).get('description', '')
            
            # Add paragraph
            if i == 1:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            p.text = f"• {title_hs}"
            p.level = 0
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = RGBColor(50, 50, 50)
            p.space_after = Pt(6)
            
            # Add description if present
            if description_hs:
                p_desc = text_frame.add_paragraph()
                p_desc.text = description_hs
                p_desc.level = 1
                p_desc.font.size = Pt(14)
                p_desc.font.color.rgb = RGBColor(100, 100, 100)
                p_desc.space_after = Pt(12)
        
        # Add more hotspots note if needed
        if len(hotspots) > 15:
            p_more = text_frame.add_paragraph()
            p_more.text = f"... and {len(hotspots) - 15} more hotspots"
            p_more.font.size = Pt(12)
            p_more.font.italic = True
            p_more.font.color.rgb = RGBColor(150, 150, 150)
    
    def _add_metadata_slide(self, prs, metadata: Dict):
        """Add slide with file metadata"""
        
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
        title_frame = title_box.text_frame
        title_frame.text = "File Information"
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(102, 126, 234)
        
        # Add metadata as bullet list
        content_box = slide.shapes.add_textbox(Inches(1), Inches(1.2), Inches(8), Inches(5.5))
        text_frame = content_box.text_frame
        text_frame.word_wrap = True
        
        for i, (key, value) in enumerate(metadata.items()):
            if key not in ['title', 'description']:
                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                
                p.text = f"{str(key).replace('_', ' ').title()}: {str(value)}"
                p.level = 0
                p.font.size = Pt(16)
                p.font.color.rgb = RGBColor(50, 50, 50)
                p.space_after = Pt(14)
